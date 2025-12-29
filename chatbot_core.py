"""
Eastern Services RAG Chatbot - Core Functions
Converted from Colab notebook to production Python code
"""

import os
import re
import json
import time
import uuid
import requests
from io import BytesIO
from datetime import datetime
from typing import List, Dict, Optional, Tuple

# Database imports
import redis
from sqlalchemy import create_engine, Column, Integer, String, Text, DateTime, text
from sqlalchemy.orm import declarative_base, sessionmaker

# Document processing imports
import docx
import PyPDF2
from pptx import Presentation
import openpyxl

# LangChain imports
from langchain_groq import ChatGroq
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough

# Cohere for embeddings
import cohere

# ============================================================================
# GLOBAL VARIABLES
# ============================================================================

# Database connections (initialized in init_databases)
redis_client = None
engine = None
SessionLocal = None
ConversationHistory = None

# Cohere client (initialized in init_embeddings)
cohere_client = None

# LLM (initialized in init_llm)
llm = None

# Data (initialized in load_documents)
LOADED_CHUNKS = []

# Session ID
SESSION_ID = str(uuid.uuid4())[:8]

# Evaluation log
evaluation_log = []

# ============================================================================
# CONFIGURATION
# ============================================================================

GITHUB_REPO_URL = "https://raw.githubusercontent.com/abdullahnaseem4802-dot/Eastern-Chatbot-Data/main/"

DATA_FILES = [
    "Ant_Control_Services.docx",
    "Bed_Bug_Treatment_Guide.docx",
    "Booking_and_Scheduling.docx",
    "Chemical_Safety_Information.docx",
    "Cockroach_Control_Services.docx",
    "Commercial_Pest_Control.docx",
    "Comprehensive_FAQ.docx",
    "Contact_Support_Information.docx",
    "Corporate_Packages_2025.docx",
    "Customer_Testimonials.docx",
    "Detailed_Pricing_2025.docx",
    "Emergency_Pest_Services.docx",
    "Fumigation_Services.docx",
    "Mosquito_Control_Services.docx",
    "Pre_Construction_Treatment.docx",
    "Rodent_Control_Solutions.docx",
    "Seasonal_Pest_Guide.docx",
    "Service_Areas_Coverage.docx",
    "Seasonal_Pest_Calendar.xlsx",
    "Pest_Identification_Chart.pptx",
    "Monthly_Service_Plans.xlsx",
    "Integrated_Pest_Management.pdf",
    "Installation_Methods_Guide.pdf",
    "DIY_Prevention_Tips.pptx",
    "Termite_Treatment_Complete.docx",
    "Warranty_Guarantee_Policy.docx"
]

# ============================================================================
# DATABASE INITIALIZATION
# ============================================================================

def init_databases(groq_key, cohere_key, supabase_uri, redis_uri):
    """Initialize Redis and Supabase connections."""
    global redis_client, engine, SessionLocal, ConversationHistory
    
    print("🔌 Connecting to databases...")
    
    # Redis connection
    try:
        redis_uri_clean = redis_uri.replace('REDIS_URI=', '').strip()
        redis_client = redis.from_url(
            redis_uri_clean,
            decode_responses=True,
            socket_connect_timeout=5
        )
        redis_client.ping()
        print("✅ Redis connected!")
    except Exception as e:
        print(f"⚠️  Redis connection failed: {e}")
        redis_client = None
    
    # Supabase connection
    try:
        engine = create_engine(
            supabase_uri,
            pool_pre_ping=True,
            pool_recycle=3600,
            pool_size=5,
            max_overflow=10,
            connect_args={
                'connect_timeout': 15,
                'application_name': 'eastern_chatbot_production'
            }
        )
        
        # Test connection
        with engine.connect() as conn:
            result = conn.execute(text("SELECT version();"))
            print("✅ Supabase connected!")
            
            # Enable PGVector if not already enabled
            result = conn.execute(text("SELECT * FROM pg_extension WHERE extname = 'vector';"))
            if not result.fetchone():
                conn.execute(text("CREATE EXTENSION IF NOT EXISTS vector;"))
                conn.commit()
                print("✅ PGVector extension enabled!")
        
        # Define schema
        Base = declarative_base()
        
        class ConversationHistory(Base):
            __tablename__ = 'conversation_history'
            
            id = Column(Integer, primary_key=True, autoincrement=True)
            session_id = Column(String(100), index=True)
            user_question = Column(Text)
            bot_response = Column(Text)
            language = Column(String(50))
            timestamp = Column(DateTime, default=datetime.utcnow)
            meta_data = Column(Text)
        
        # Create tables
        Base.metadata.create_all(engine)
        SessionLocal = sessionmaker(bind=engine)
        print("✅ Database tables ready!")
        
    except Exception as e:
        print(f"❌ Supabase connection failed: {e}")
        engine = None
        SessionLocal = None
    
    return redis_client is not None and engine is not None

# ============================================================================
# DOCUMENT LOADING
# ============================================================================

def download_file_from_github(url, filename):
    """Download a file from GitHub."""
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        return response.content
    except Exception as e:
        print(f"   ⚠️  Failed to download {filename}: {e}")
        return None

def extract_text_from_docx(file_content):
    """Extract text from .docx file."""
    try:
        doc = docx.Document(BytesIO(file_content))
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    full_text.append(row_text)
        
        return '\n\n'.join(full_text)
    except Exception as e:
        print(f"   ❌ Error extracting DOCX: {e}")
        return None

def extract_text_from_pdf(file_content):
    """Extract text from PDF file."""
    try:
        pdf_file = BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        full_text = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text.strip():
                full_text.append(text)
        
        return '\n\n'.join(full_text)
    except Exception as e:
        print(f"   ❌ Error extracting PDF: {e}")
        return None

def extract_text_from_pptx(file_content):
    """Extract text from PowerPoint file."""
    try:
        pptx_file = BytesIO(file_content)
        prs = Presentation(pptx_file)
        
        full_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                full_text.append(f"Slide {slide_num}:\n" + '\n'.join(slide_text))
        
        return '\n\n'.join(full_text)
    except Exception as e:
        print(f"   ❌ Error extracting PPTX: {e}")
        return None

def extract_text_from_excel(file_content):
    """Extract text from Excel file."""
    try:
        excel_file = BytesIO(file_content)
        workbook = openpyxl.load_workbook(excel_file, data_only=True)
        
        full_text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = [f"Sheet: {sheet_name}"]
            
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    sheet_text.append(row_text)
            
            if len(sheet_text) > 1:
                full_text.append('\n'.join(sheet_text))
        
        return '\n\n'.join(full_text)
    except Exception as e:
        print(f"   ❌ Error extracting Excel: {e}")
        return None

def extract_text_from_file(file_content, filename):
    """Extract text based on file extension."""
    extension = filename.lower().split('.')[-1]
    
    extractors = {
        'docx': extract_text_from_docx,
        'pdf': extract_text_from_pdf,
        'pptx': extract_text_from_pptx,
        'xlsx': extract_text_from_excel,
        'xls': extract_text_from_excel
    }
    
    extractor = extractors.get(extension)
    if extractor:
        return extractor(file_content)
    else:
        print(f"   ⚠️  Unsupported file type: .{extension}")
        return None

def simple_chunk_text(text, chunk_size=1000, overlap=200):
    """Simple text chunking."""
    chunks = []
    start = 0
    text_length = len(text)
    
    while start < text_length:
        end = start + chunk_size
        chunk = text[start:end]
        
        if end < text_length:
            for separator in ['\n\n', '\n', '. ', ' ']:
                last_sep = chunk.rfind(separator)
                if last_sep > chunk_size * 0.5:
                    chunk = text[start:start + last_sep + len(separator)]
                    break
        
        chunks.append(chunk.strip())
        start = end - overlap
    
    return chunks

def load_documents_from_github(repo_url, file_list):
    """Load all documents from GitHub repository."""
    all_documents = []
    print(f"📥 Loading documents from GitHub...")
    print(f"   Total files: {len(file_list)}")
    
    for filename in file_list:
        file_url = repo_url + filename
        print(f"   Downloading: {filename}")
        
        file_content = download_file_from_github(file_url, filename)
        
        if file_content:
            try:
                text = extract_text_from_file(file_content, filename)
                
                if text:
                    all_documents.append({
                        'content': text,
                        'source': filename,
                        'file_type': filename.split('.')[-1]
                    })
                    print(f"   ✅ Loaded: {filename}")
            except Exception as e:
                print(f"   ❌ Error loading {filename}: {e}")
    
    print(f"✅ Loaded {len(all_documents)} documents")
    return all_documents

def create_chunks_from_documents(documents):
    """Create chunks from all documents."""
    print("📄 Creating chunks...")
    all_chunks = []
    
    for doc in documents:
        text = doc['content']
        source = doc['source']
        
        text_chunks = simple_chunk_text(text, chunk_size=1000, overlap=200)
        
        for i, chunk_text in enumerate(text_chunks):
            chunk = {
                'text': chunk_text,
                'metadata': {
                    'source': source,
                    'file_type': doc['file_type'],
                    'chunk_index': i,
                    'total_chunks': len(text_chunks)
                }
            }
            all_chunks.append(chunk)
    
    print(f"✅ Created {len(all_chunks)} chunks")
    return all_chunks

# ============================================================================
# EMBEDDINGS AND RETRIEVAL
# ============================================================================

def init_embeddings(cohere_key):
    """Initialize Cohere embeddings."""
    global cohere_client
    
    print("🔌 Initializing Cohere embeddings...")
    try:
        cohere_client = cohere.Client(cohere_key)
        print("✅ Cohere client initialized!")
        return True
    except Exception as e:
        print(f"❌ Cohere initialization failed: {e}")
        return False

def setup_pgvector_table():
    """Create embeddings table in Supabase."""
    print("📋 Setting up PGVector table...")
    
    try:
        with engine.connect() as conn:
            # Create embeddings table
            conn.execute(text("""
                CREATE TABLE IF NOT EXISTS document_embeddings (
                    id SERIAL PRIMARY KEY,
                    chunk_text TEXT NOT NULL,
                    embedding vector(1024),
                    metadata JSONB,
                    created_at TIMESTAMP DEFAULT NOW()
                );
            """))
            
            # Create index
            conn.execute(text("""
                CREATE INDEX IF NOT EXISTS embedding_idx
                ON document_embeddings
                USING ivfflat (embedding vector_cosine_ops)
                WITH (lists = 100);
            """))
            
            conn.commit()
            print("✅ PGVector table ready!")
            
            # Check existing embeddings
            result = conn.execute(text("SELECT COUNT(*) FROM document_embeddings;"))
            existing_count = result.fetchone()[0]
            
            return existing_count
    except Exception as e:
        print(f"❌ Error setting up PGVector table: {e}")
        return 0

def generate_and_store_embeddings(chunks):
    """Generate and store embeddings for all chunks."""
    print(f"🔄 Generating embeddings for {len(chunks)} chunks...")
    
    try:
        batch_size = 90
        total_embedded = 0
        
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i + batch_size]
            batch_texts = [chunk['text'] for chunk in batch]
            
            print(f"   Processing batch {i//batch_size + 1}/{(len(chunks)-1)//batch_size + 1}...")
            
            # Get embeddings from Cohere
            response = cohere_client.embed(
                texts=batch_texts,
                model='embed-english-v3.0',
                input_type='search_document'
            )
            
            embeddings_batch = response.embeddings
            
            # Insert into Supabase
            raw_conn = engine.raw_connection()
            try:
                cursor = raw_conn.cursor()
                
                for chunk, embedding in zip(batch, embeddings_batch):
                    embedding_str = '[' + ','.join(map(str, embedding)) + ']'
                    metadata_json = json.dumps(chunk['metadata'])
                    
                    cursor.execute("""
                        INSERT INTO document_embeddings (chunk_text, embedding, metadata)
                        VALUES (%s, %s::vector, %s::jsonb)
                    """, (chunk['text'], embedding_str, metadata_json))
                    
                    total_embedded += 1
                
                raw_conn.commit()
                cursor.close()
            finally:
                raw_conn.close()
            
            print(f"   ✅ Embedded {total_embedded}/{len(chunks)} chunks")
        
        print(f"✅ All embeddings stored!")
        return True
    except Exception as e:
        print(f"❌ Error generating embeddings: {e}")
        return False

# Urdu transliteration dictionary (200+ words)
URDU_TO_ROMAN = {
    # Core pest terms
    'دیمک': 'deemak termite white ant wood damage',
    'کھٹمل': 'khatmal bed bug bedbug mattress',
    'کاکروچ': 'cockroach roach kitchen pest',
    'مچھر': 'mosquito machar dengue malaria',
    'چوہے': 'rat rodent mouse chooha chuhay',
    'چوہا': 'rat rodent mouse chooha',
    'چیونٹی': 'ant chiti sugar ant',
    
    # Service terms
    'علاج': 'ilaj treatment control service',
    'سروس': 'service treatment ilaj',
    'خدمات': 'services khidmat',
    'کنٹرول': 'control treatment',
    'سپرے': 'spray chemical treatment',
    'فیومیگیشن': 'fumigation',
    'معائنہ': 'inspection check survey',
    
    # Pricing terms
    'ریٹ': 'rate price cost qeemat',
    'قیمت': 'qeemat price rate cost',
    'رقم': 'amount money price',
    'روپے': 'rupees rs money',
    'فی': 'per',
    'مربع': 'square sq',
    'فٹ': 'foot ft feet',
    'مرلہ': 'marla',
    'کنال': 'kanal',
    
    # Guarantee terms
    'گارنٹی': 'guarantee warranty assurance',
    'ضمانت': 'guarantee warranty',
    'وارنٹی': 'warranty guarantee',
    
    # Time terms
    'کتنی': 'kitni how much how many',
    'کتنا': 'kitna how much',
    'کب': 'kab when',
    'کیسے': 'kaisay how',
    'کیوں': 'kyun why',
    'وقت': 'waqt time duration',
    'دن': 'day din',
    'ماہ': 'month mahina',
    'سال': 'year saal',
    
    # Location terms
    'گھر': 'ghar home house residence',
    'مکان': 'house home makan',
    'علاقہ': 'area ilaqa region',
    'علاقے': 'areas ilaqay regions',
    'شہر': 'city shahar',
    'لکڑی': 'wood lakri timber',
    'تعمیر': 'construction building tameer',
    
    # Question words
    'کیا': 'kya kia what',
    'کون': 'kon kaun which who what',
    'کہاں': 'kahan where',
    'کس': 'kis which',
    'کونسا': 'konsa which',
    'کونسی': 'konsi which',
    'کونسے': 'konsay which',
    
    # Common verbs
    'ہے': 'hai is',
    'ہیں': 'hain are',
    'ہو': 'ho',
    'گا': 'ga will',
    'گی': 'gi will',
    'گے': 'gay will',
    'ہوگا': 'hoga will be',
    'ہوگی': 'hogi will be',
    'ہوگے': 'hogay will be',
    'کریں': 'karein do',
    'کرنا': 'karna to do',
    'دیں': 'dein give',
    'دینا': 'dena to give',
    'آئیں': 'aayein come',
    'آنا': 'aana to come',
    'جائیں': 'jayein go',
    'جانا': 'jana to go',
    
    # Possessives
    'میرے': 'meray mery my',
    'میری': 'meri my',
    'میرا': 'mera my',
    'آپ': 'aap you',
    'آپکا': 'aapka your',
    'آپکی': 'aapki your',
    'آپکے': 'aapkay your',
    'ہمارے': 'humare our',
    'ہماری': 'humari our',
    'ہمارا': 'humara our',
    
    # Prepositions
    'میں': 'mein in',
    'کے': 'ke kay of',
    'کی': 'ki of',
    'کا': 'ka of',
    'سے': 'se say from with',
    'کو': 'ko to',
    'نے': 'ne',
    'پر': 'par on',
    'لیے': 'liye liyay for',
    'بارے': 'baray about',
    'ساتھ': 'saath with',
    
    # Adjectives
    'بہت': 'bohat bahut very much',
    'زیادہ': 'zyada more',
    'کم': 'kam less',
    'بڑا': 'bara big large',
    'چھوٹا': 'chota small',
    'اچھا': 'acha good',
    'بہترین': 'behtareen best excellent',
    'پیشہ ورانہ': 'professional peshawrana',
    'معیاری': 'quality standard',
    
    # Common nouns
    'مسئلہ': 'masla problem issue',
    'حل': 'solution hal',
    'طریقہ': 'method tariqa',
    'پیکج': 'package',
    'بکنگ': 'booking',
    'شیڈولنگ': 'scheduling',
    'رابطہ': 'contact rabta',
    'مدد': 'help madad',
    'معلومات': 'information maloomat',
    'تفصیل': 'details tafseel',
    
    # Negation
    'نہیں': 'nahi no not',
    'نہ': 'na no not',
    
    # Conjunctions
    'اور': 'aur and',
    'یا': 'ya or',
    'لیکن': 'lekin but',
    'مگر': 'magar but',
    'تو': 'to then',
    'اگر': 'agar if',
    
    # Polite terms
    'براہ': 'barah please',
    'کرم': 'karam kindness',
    'شکریہ': 'shukriya thanks',
    'معذرت': 'maazrat sorry',
    'خوش': 'khush happy',
    'آمدید': 'aamdeed welcome',
    
    # Business terms
    'کمپنی': 'company',
    'فرم': 'firm',
    'ٹیم': 'team',
    'ماہر': 'expert mahir',
    'تجربہ': 'experience tajurba',
    'ضرورت': 'need zaroorat',
    'فوری': 'urgent fori',
    'ایمرجنسی': 'emergency',
    
    # Additional context
    'پہلے': 'pehlay before first',
    'بعد': 'baad after',
    'دوران': 'during doran',
    'ختم': 'khatam end finish',
    'شروع': 'start shuru',
    'مکمل': 'complete mukammal',
}

def transliterate_urdu_to_roman(text):
    """Comprehensive Urdu to Roman transliteration."""
    urdu_pattern = r'[\u0600-\u06FF]'
    if not re.search(urdu_pattern, text):
        return text
    
    roman_text = text
    for urdu_word, roman_word in URDU_TO_ROMAN.items():
        roman_text = roman_text.replace(urdu_word, roman_word)
    
    # Remove remaining non-Roman characters
    roman_text = re.sub(r'[^\w\s]', ' ', roman_text)
    roman_text = re.sub(r'\s+', ' ', roman_text)
    
    return roman_text.strip()

def retrieve_similar_chunks(query, k=3):
    """Retrieve similar chunks from PGVector with keyword re-ranking."""
    try:
        # Transliterate Urdu to Roman
        search_query = transliterate_urdu_to_roman(query)
        
        # Generate query embedding
        query_response = cohere_client.embed(
            texts=[search_query],
            model='embed-english-v3.0',
            input_type='search_query'
        )
        query_embedding = query_response.embeddings[0]
        query_embedding_str = '[' + ','.join(map(str, query_embedding)) + ']'
        
        # Retrieve more chunks for re-ranking
        retrieval_k = k + 2
        
        # Search for similar chunks
        raw_conn = engine.raw_connection()
        try:
            cursor = raw_conn.cursor()
            
            cursor.execute("""
                SELECT chunk_text, metadata,
                       1 - (embedding <=> %s::vector) as similarity
                FROM document_embeddings
                ORDER BY embedding <=> %s::vector
                LIMIT %s
            """, (query_embedding_str, query_embedding_str, retrieval_k))
            
            rows = cursor.fetchall()
            cursor.close()
            
            documents = []
            for row in rows:
                metadata = row[1] if isinstance(row[1], dict) else (json.loads(row[1]) if row[1] else {})
                doc = {
                    'page_content': row[0],
                    'metadata': metadata,
                    'similarity': row[2]
                }
                documents.append(doc)
            
            # Keyword-based re-ranking
            query_lower = search_query.lower()
            
            pest_keywords = {
                'deemak': ['termite', 'deemak', 'white ant', 'wood'],
                'termite': ['termite', 'deemak', 'white ant', 'wood'],
                'khatmal': ['bed bug', 'khatmal', 'bedbug', 'mattress'],
                'bed': ['bed bug', 'khatmal', 'bedbug', 'mattress'],
                'bug': ['bed bug', 'khatmal', 'bedbug', 'mattress'],
                'cockroach': ['cockroach', 'roach', 'kitchen'],
                'mosquito': ['mosquito', 'malaria', 'dengue', 'machar'],
                'rat': ['rat', 'rodent', 'mouse', 'chuhay', 'chooha'],
                'chuhay': ['rat', 'rodent', 'mouse', 'chuhay', 'chooha'],
                'ant': ['ant', 'chiti', 'sugar ant'],
                'guarantee': ['guarantee', 'warranty', 'guaranteed', 'assurance'],
                'rate': ['rate', 'price', 'cost', 'pricing', 'charges', 'rupees', 'rs', 'qeemat'],
                'service': ['service', 'services', 'provide', 'offer', 'ilaj', 'treatment'],
                'marla': ['marla', 'sq ft', 'square', 'area', 'size'],
                'kanal': ['kanal', 'sq ft', 'square', 'area', 'size'],
            }
            
            query_relevant_keywords = []
            for key, keywords in pest_keywords.items():
                if any(kw in query_lower for kw in [key]):
                    query_relevant_keywords.extend(keywords)
            
            query_relevant_keywords = list(set(query_relevant_keywords))
            
            # Re-rank documents
            if query_relevant_keywords:
                for doc in documents:
                    doc_text_lower = doc['page_content'].lower()
                    keyword_matches = sum(1 for kw in query_relevant_keywords if kw in doc_text_lower)
                    
                    if keyword_matches > 0:
                        boost = min(0.5, keyword_matches * 0.2)
                        doc['similarity'] = min(1.0, doc['similarity'] + boost)
                        doc['keyword_matches'] = keyword_matches
                    else:
                        doc['similarity'] = doc['similarity'] * 0.8
                        doc['keyword_matches'] = 0
                
                documents.sort(key=lambda x: x['similarity'], reverse=True)
            
            return documents[:k]
        
        finally:
            raw_conn.close()
    
    except Exception as e:
        print(f"❌ Error retrieving chunks: {e}")
        return []

# ============================================================================
# LLM INITIALIZATION
# ============================================================================

def init_llm(groq_key):
    """Initialize Groq LLM."""
    global llm
    
    print("🔌 Initializing Groq LLM...")
    try:
        llm = ChatGroq(
            api_key=groq_key,
            model="llama-3.3-70b-versatile",
            temperature=0.0
        )
        print("✅ Groq LLM initialized (Llama 3.3 70B)!")
        return True
    except Exception as e:
        print(f"❌ LLM initialization failed: {e}")
        return False

# ============================================================================
# LANGUAGE DETECTION
# ============================================================================

def detect_greeting(text):
    """Detect if text contains a greeting."""
    text_lower = text.lower()
    urdu_pattern = r'[\u0600-\u06FF]'
    has_urdu_script = re.search(urdu_pattern, text)
    
    if 'السلام عليكم' in text or 'السلام' in text:
        return 'urdu_greeting'
    
    if 'assalam' in text_lower or 'assalamu' in text_lower:
        if has_urdu_script:
            return 'urdu_greeting'
        return 'roman_urdu_greeting'
    
    if any(greet in text_lower for greet in ['hello', 'hi', 'hey', 'greetings']):
        return 'english_greeting'
    
    return None

def detect_language(text):
    """Detect language (English, Roman Urdu, or Pure Urdu)."""
    text_lower = text.lower()
    
    # Check for Urdu script
    urdu_pattern = r'[\u0600-\u06FF]'
    if re.search(urdu_pattern, text):
        return 'urdu'
    
    # Punjabi keywords (mapped to Roman Urdu)
    punjabi_keywords = [
        'tuhada', 'tussin', 'tusi', 'ah pai', 'ki hall', 'lag gi',
        'kaar cha', 'vich', 'da', 'di', 'oday', 'tay', 'karwana',
        'chana', 'karday', 'keri', 'oda', 'nu', 'te', 'ae',
        'samajh layya', 'dasso', 'layi', 'nal', 'baray', 'chahoge',
        'meray kaar', 'wood da kam', 'treatement da', 'ki ah'
    ]
    punjabi_count = sum(1 for keyword in punjabi_keywords if keyword in text_lower)
    
    # Roman Urdu keywords
    urdu_keywords = [
        'aapka', 'aapko', 'samajh gaya', 'hain', 'kia', 'hai',
        'karna', 'hoga', 'meray', 'ghar', 'deemak', 'rate',
        'ho ga', 'mein', 'ke', 'ka', 'ki', 'kitni', 'kitna',
        'chahta', 'chahiye', 'zaroorat', 'problem', 'masla',
        'ilaj', 'khatam', 'control', 'service', 'guarantee'
    ]
    urdu_count = sum(1 for keyword in urdu_keywords if keyword in text_lower)
    
    # English keywords
    english_keywords = [
        'what', 'which', 'how', 'when', 'where', 'why',
        'is', 'are', 'was', 'were', 'has', 'have', 'had',
        'the', 'this', 'that', 'these', 'those',
        'treatment', 'process', 'guarantee', 'period',
        'clients', 'worked', 'with', 'services', 'offer'
    ]
    english_count = sum(1 for keyword in english_keywords if keyword in text_lower)
    
    # Map Punjabi to Roman Urdu
    if punjabi_count >= 2:
        return 'roman_urdu'
    
    if urdu_count >= 3 and urdu_count > english_count:
        return 'roman_urdu'
    
    if english_count >= 3:
        return 'english'
    
    if any(word in text_lower for word in ['kitni', 'kitna', 'kia', 'kya']) and english_count < 2:
        return 'roman_urdu'
    
    if any(word in text_lower for word in ['tusi', 'tuhada', 'karday', 'keri', 'vich', 'da di']):
        return 'roman_urdu'
    
    return 'english'

def clean_response(response, language, original_question=None):
    """Remove unwanted meta-commentary and clean up response."""
    patterns = [
        r'Q:\s*.*?\n\s*A:\s*',
        r'Based on the provided context.*?:\s*',
        r'According to the context.*?:\s*',
        r'RULE \d+.*?:\s*',
        r'INSTRUCTIONS?:.*?:\s*',
        r'Since the question.*?:\s*',
        r'The question is in.*?:\s*',
        r'LANGUAGE DETECTION.*?:\s*',
        r'I will.*?:\s*',
        r'Let me.*?:\s*',
        r'\([^)]*\)',
        r'Assalam-o-Alaikum!?\s*',
        r'السلام عليكم\s*',
    ]
    
    for pattern in patterns:
        response = re.sub(pattern, '', response, flags=re.IGNORECASE | re.DOTALL)
    
    # Remove repeated question if present
    if original_question:
        question_clean = original_question.strip()
        question_clean = re.sub(r'Assalam-o-Alaikum!?\s*', '', question_clean, flags=re.IGNORECASE)
        question_clean = re.sub(r'السلام عليكم\s*', '', question_clean)
        question_clean = question_clean.strip()
        
        response_stripped = response.strip()
        
        if question_clean and response_stripped.startswith(question_clean):
            temp_response = response_stripped[len(question_clean):].strip()
            if len(temp_response) > 20:
                response = temp_response
    
    response = response.strip()
    response = re.sub(r'\n\n+', '\n', response)
    response = re.sub(r'\s+', ' ', response)
    
    return response

# ============================================================================
# RAG PROMPTS (continued in next part due to length)
# ============================================================================
